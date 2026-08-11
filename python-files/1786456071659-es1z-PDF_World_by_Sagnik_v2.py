# PDF World by Sagnik — v2.0
# A local, colourful PDF toolkit for Windows.
# No LibreOffice is required for the built-in converters.
#
# Built-in converters:
#   To PDF: TXT, HTML/HTM, DOCX, XLSX, PPTX, PNG/JPG/JPEG/BMP/GIF/TIFF/WEBP
#   From PDF: PNG/JPG, TXT, HTML, DOCX, XLSX, PPTX
#
# PDF tools:
#   Merge, Split/Extract, Rotate, Delete Pages, Reorder Pages,
#   Compress, Watermark, Password Protect, Unlock (known password)
#
# Important:
#   Old binary Office files (.DOC/.XLS/.PPT) cannot be faithfully converted
#   without an Office-compatible engine. If Microsoft Office is installed,
#   an optional Windows/Office bridge can be added later.
#
# Created for: PDF World by Sagnik

import os
import sys
import html
import re
import math
import threading
import traceback
from pathlib import Path

import tkinter as tk
from tkinter import ttk, filedialog, messagebox, simpledialog

# Optional drag/drop. The program still works without it.
try:
    from tkinterdnd2 import DND_FILES, TkinterDnD
    DND_OK = True
except Exception:
    DND_OK = False

# Required packages.
MISSING = []
try:
    import fitz  # PyMuPDF
except Exception:
    MISSING.append("PyMuPDF")

try:
    from PIL import Image, ImageOps, ImageTk
except Exception:
    MISSING.append("Pillow")

try:
    from pypdf import PdfReader, PdfWriter, PdfMerger
except Exception:
    MISSING.append("pypdf")

try:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.lib.pagesizes import A4, LETTER, landscape
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        PageBreak, Image as RLImage, KeepTogether
    )
except Exception:
    MISSING.append("reportlab")

try:
    from docx import Document
    from docx.shared import Inches, Pt
except Exception:
    MISSING.append("python-docx")

try:
    from openpyxl import load_workbook, Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
except Exception:
    MISSING.append("openpyxl")

try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches
except Exception:
    MISSING.append("python-pptx")


if MISSING:
    raise SystemExit(
        "PDF World is missing these packages:\n\n  "
        + "\n  ".join(MISSING)
        + "\n\nInstall them with:\n"
        "py -m pip install -r requirements_pdf_world_v2.txt"
    )


APP_NAME = "PDF World by Sagnik"
VERSION = "2.0"

# ---------- Theme ----------
BG = "#F4F7FB"
WHITE = "#FFFFFF"
TEXT = "#172033"
MUTED = "#6B7280"
BORDER = "#E5E7EB"
DARK = "#111827"

PURPLE = "#6D28D9"
PURPLE2 = "#8B5CF6"
BLUE = "#2563EB"
CYAN = "#0891B2"
GREEN = "#16A34A"
ORANGE = "#EA580C"
PINK = "#DB2777"
RED = "#DC2626"
TEAL = "#0F766E"

CARD_COLORS = [PURPLE, BLUE, CYAN, GREEN, ORANGE, PINK, TEAL, RED]


# ---------- Utility ----------
def safe_text(value):
    """Convert arbitrary content to safe plain text."""
    if value is None:
        return ""
    return str(value).replace("\x00", "")


def esc_pdf_text(value):
    """Escape text for ReportLab Paragraph."""
    return html.escape(safe_text(value), quote=False)


def unique_path(path: Path) -> Path:
    """Avoid overwriting existing files."""
    if not path.exists():
        return path
    stem, suffix = path.stem, path.suffix
    for i in range(2, 10000):
        candidate = path.with_name(f"{stem} ({i}){suffix}")
        if not candidate.exists():
            return candidate
    raise RuntimeError(f"Could not create a unique filename for {path.name}")


def choose_output_dir():
    return filedialog.askdirectory(title="Choose output folder")


def choose_pdf():
    return filedialog.askopenfilename(
        title="Choose a PDF",
        filetypes=[("PDF files", "*.pdf")]
    )


def choose_save(ext, description, initial):
    return filedialog.asksaveasfilename(
        title="Save output",
        defaultextension=ext,
        initialfile=initial,
        filetypes=[(description, "*" + ext)]
    )


def read_text_file(path: Path) -> str:
    for enc in ("utf-8-sig", "utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def wrap_long_text(text, width=110):
    lines = []
    for raw in safe_text(text).splitlines():
        if len(raw) <= width:
            lines.append(raw)
        else:
            for i in range(0, len(raw), width):
                lines.append(raw[i:i + width])
    return lines


# ---------- PDF creation ----------
class PDFBuilder:
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self.body = ParagraphStyle(
            "PWBody",
            parent=self.styles["BodyText"],
            fontName="Helvetica",
            fontSize=9.5,
            leading=13,
            spaceAfter=5,
        )
        self.heading = ParagraphStyle(
            "PWHeading",
            parent=self.styles["Heading2"],
            fontName="Helvetica-Bold",
            fontSize=14,
            leading=18,
            spaceBefore=6,
            spaceAfter=8,
        )
        self.small = ParagraphStyle(
            "PWSmall",
            parent=self.body,
            fontSize=8,
            leading=10,
        )
        self.title = ParagraphStyle(
            "PWTitle",
            parent=self.styles["Title"],
            fontName="Helvetica-Bold",
            fontSize=22,
            leading=26,
            alignment=TA_CENTER,
            spaceAfter=14,
        )

    def document(self, dest, page_size=A4):
        return SimpleDocTemplate(
            str(dest),
            pagesize=page_size,
            rightMargin=16 * mm,
            leftMargin=16 * mm,
            topMargin=16 * mm,
            bottomMargin=16 * mm,
            title="PDF World by Sagnik",
            author="PDF World by Sagnik",
        )

    def txt_to_pdf(self, src: Path, dest: Path):
        text = read_text_file(src)
        story = []
        # IMPORTANT: do NOT add src.name as a title.
        for line in text.splitlines():
            if line.strip():
                story.append(Paragraph(esc_pdf_text(line), self.body))
            else:
                story.append(Spacer(1, 5))
        if not story:
            story = [Paragraph("", self.body)]
        self.document(dest).build(story)

    def html_to_pdf(self, src: Path, dest: Path):
        raw = read_text_file(src)

        # Remove non-visible content.
        raw = re.sub(r"<script\b[^>]*>.*?</script>", "", raw,
                     flags=re.I | re.S)
        raw = re.sub(r"<style\b[^>]*>.*?</style>", "", raw,
                     flags=re.I | re.S)

        # Preserve common structure rather than dumping the filename.
        raw = re.sub(r"<br\s*/?>", "\n", raw, flags=re.I)
        raw = re.sub(r"</p\s*>", "\n\n", raw, flags=re.I)
        raw = re.sub(r"</div\s*>", "\n", raw, flags=re.I)
        raw = re.sub(r"<li\b[^>]*>", "• ", raw, flags=re.I)
        raw = re.sub(r"</li\s*>", "\n", raw, flags=re.I)
        raw = re.sub(r"</h[1-6]\s*>", "\n", raw, flags=re.I)
        text = re.sub(r"<[^>]+>", "", raw)
        text = html.unescape(text)
        text = re.sub(r"\n{3,}", "\n\n", text)

        story = []
        for line in text.splitlines():
            line = line.strip()
            if not line:
                story.append(Spacer(1, 5))
            else:
                story.append(Paragraph(esc_pdf_text(line), self.body))
        if not story:
            story = [Paragraph("", self.body)]
        self.document(dest).build(story)

    def docx_to_pdf(self, src: Path, dest: Path):
        docx = Document(str(src))
        story = []

        for p in docx.paragraphs:
            text = safe_text(p.text)
            if not text.strip():
                story.append(Spacer(1, 6))
                continue

            style_name = safe_text(getattr(p.style, "name", "")).lower()
            if "heading 1" in style_name:
                style = self.heading
            elif "heading" in style_name:
                style = self.heading
            else:
                style = self.body

            # Basic bold/italic runs are converted into simple Paragraph markup.
            pieces = []
            for run in p.runs:
                t = esc_pdf_text(run.text)
                if not t:
                    continue
                if run.bold and run.italic:
                    t = f"<b><i>{t}</i></b>"
                elif run.bold:
                    t = f"<b>{t}</b>"
                elif run.italic:
                    t = f"<i>{t}</i>"
                pieces.append(t)

            paragraph = "".join(pieces) or esc_pdf_text(text)
            story.append(Paragraph(paragraph, style))

        for table in docx.tables:
            data = []
            for row in table.rows:
                data.append([esc_pdf_text(cell.text) for cell in row.cells])

            if data:
                col_count = max(len(r) for r in data)
                data = [r + [""] * (col_count - len(r)) for r in data]
                t = Table(data, repeatRows=1, hAlign="LEFT")
                t.setStyle(TableStyle([
                    ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CBD5E1")),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EDE9FE")),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                    ("FONTSIZE", (0, 0), (-1, -1), 7.5),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 4),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ]))
                story.extend([Spacer(1, 8), t, Spacer(1, 10)])

        if not story:
            story = [Paragraph("", self.body)]
        self.document(dest).build(story)

    def xlsx_to_pdf(self, src: Path, dest: Path):
        wb = load_workbook(str(src), data_only=True, read_only=True)
        story = []

        for index, ws in enumerate(wb.worksheets):
            if index:
                story.append(PageBreak())
            story.append(Paragraph(esc_pdf_text(ws.title), self.heading))

            rows = []
            for row in ws.iter_rows(values_only=True):
                vals = [safe_text(v) for v in row]
                while vals and vals[-1] == "":
                    vals.pop()
                if vals:
                    rows.append(vals)

            if not rows:
                story.append(Paragraph("(Empty worksheet)", self.body))
                continue

            # Keep table width manageable. All columns are retained up to a safe limit.
            max_cols = min(max(len(r) for r in rows), 18)
            rows = [
                r[:max_cols] + [""] * max(0, max_cols - len(r))
                for r in rows
            ]

            table = Table(rows, repeatRows=1, hAlign="LEFT")
            table.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#CBD5E1")),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#DBEAFE")),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 6.5),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 3),
                ("RIGHTPADDING", (0, 0), (-1, -1), 3),
            ]))
            story.append(table)

        wb.close()
        self.document(dest, landscape(A4)).build(story or [Paragraph("", self.body)])

    def pptx_to_pdf(self, src: Path, dest: Path):
        prs = Presentation(str(src))
        story = []

        for i, slide in enumerate(prs.slides):
            if i:
                story.append(PageBreak())
            story.append(Paragraph(f"Slide {i + 1}", self.heading))

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = safe_text(shape.text).strip()
                if text:
                    for line in text.splitlines():
                        story.append(Paragraph(esc_pdf_text(line), self.body))

        self.document(dest, landscape(A4)).build(
            story or [Paragraph("", self.body)]
        )

    def image_to_pdf(self, src: Path, dest: Path):
        with Image.open(src) as im:
            rgb = ImageOps.exif_transpose(im).convert("RGB")
            width, height = rgb.size
            # Use ReportLab so very large images are scaled to the page instead
            # of creating an enormous PDF page.
            page_w, page_h = A4
            max_w, max_h = page_w - 32 * mm, page_h - 32 * mm
            scale = min(max_w / width, max_h / height, 1.0)
            w, h = width * scale, height * scale

            temp = dest.with_name(dest.stem + "__temp.png")
            rgb.save(temp, "PNG")
            story = [
                RLImage(str(temp), width=w, height=h),
            ]
            self.document(dest).build(story)
            try:
                temp.unlink()
            except OSError:
                pass

    def images_to_pdf(self, files, dest):
        story = []
        temp_files = []
        try:
            for src in files:
                with Image.open(src) as im:
                    rgb = ImageOps.exif_transpose(im).convert("RGB")
                    temp = dest.with_name(
                        f"__pw_tmp_{len(temp_files)}_{src.stem}.png"
                    )
                    rgb.save(temp, "PNG")
                    temp_files.append(temp)

                    page_w, page_h = A4
                    max_w, max_h = page_w - 32 * mm, page_h - 32 * mm
                    scale = min(max_w / rgb.width, max_h / rgb.height, 1.0)
                    story.append(RLImage(
                        str(temp),
                        width=rgb.width * scale,
                        height=rgb.height * scale
                    ))
                    if src != files[-1]:
                        story.append(PageBreak())

            self.document(dest).build(story or [Paragraph("", self.body)])
        finally:
            for temp in temp_files:
                try:
                    temp.unlink()
                except OSError:
                    pass


# ---------- PDF operations ----------
class PDFTools:
    @staticmethod
    def merge(files, dest):
        merger = PdfMerger()
        try:
            for f in files:
                merger.append(str(f))
            with open(dest, "wb") as out:
                merger.write(out)
        finally:
            merger.close()

    @staticmethod
    def extract(src, dest, pages):
        reader = PdfReader(str(src))
        writer = PdfWriter()
        total = len(reader.pages)

        for p in pages:
            if 1 <= p <= total:
                writer.add_page(reader.pages[p - 1])

        with open(dest, "wb") as out:
            writer.write(out)

    @staticmethod
    def rotate(src, dest, degrees):
        reader = PdfReader(str(src))
        writer = PdfWriter()
        degrees = degrees % 360
        for page in reader.pages:
            page.rotate(degrees)
            writer.add_page(page)
        with open(dest, "wb") as out:
            writer.write(out)

    @staticmethod
    def delete_pages(src, dest, delete_set):
        reader = PdfReader(str(src))
        writer = PdfWriter()
        for i, page in enumerate(reader.pages, start=1):
            if i not in delete_set:
                writer.add_page(page)
        with open(dest, "wb") as out:
            writer.write(out)

    @staticmethod
    def reorder(src, dest, order):
        reader = PdfReader(str(src))
        writer = PdfWriter()
        total = len(reader.pages)
        for p in order:
            if 1 <= p <= total:
                writer.add_page(reader.pages[p - 1])
        with open(dest, "wb") as out:
            writer.write(out)

    @staticmethod
    def compress(src, dest):
        # qpdf-like deep compression is not available in pypdf.
        # This method removes some duplicate object streams where possible.
        reader = PdfReader(str(src))
        writer = PdfWriter()
        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)

        with open(dest, "wb") as out:
            writer.write(out)

    @staticmethod
    def watermark(src, dest, text):
        from reportlab.pdfgen import canvas

        watermark_path = Path(dest).with_name("__pw_watermark.pdf")
        c = canvas.Canvas(str(watermark_path), pagesize=A4)
        width, height = A4
        c.saveState()
        c.setFillColor(colors.Color(0.35, 0.35, 0.35, alpha=0.18))
        c.setFont("Helvetica-Bold", 34)
        c.translate(width / 2, height / 2)
        c.rotate(35)
        c.drawCentredString(0, 0, text[:70])
        c.restoreState()
        c.save()

        try:
            base = PdfReader(str(src))
            mark = PdfReader(str(watermark_path))
            writer = PdfWriter()
            wm = mark.pages[0]
            for page in base.pages:
                page.merge_page(wm)
                writer.add_page(page)
            with open(dest, "wb") as out:
                writer.write(out)
        finally:
            try:
                watermark_path.unlink()
            except OSError:
                pass

    @staticmethod
    def protect(src, dest, password):
        reader = PdfReader(str(src))
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(password)
        with open(dest, "wb") as out:
            writer.write(out)

    @staticmethod
    def unlock(src, dest, password):
        reader = PdfReader(str(src))
        if reader.is_encrypted:
            if not reader.decrypt(password):
                raise ValueError("The supplied PDF password is incorrect.")
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        with open(dest, "wb") as out:
            writer.write(out)


# ---------- GUI ----------
class PDFWorldApp:
    def __init__(self, root):
        self.root = root
        self.root.title(f"{APP_NAME}  •  v{VERSION}")
        self.root.geometry("1180x760")
        self.root.minsize(980, 650)
        self.root.configure(bg=BG)

        self.status = tk.StringVar(value="Ready")
        self.progress = tk.DoubleVar(value=0)
        self.current_page = tk.StringVar(value="Home")
        self.drag_files = []

        self.builder = PDFBuilder()

        self.setup_styles()
        self.build_shell()
        self.show_home()

    # ----- Style helpers -----
    def setup_styles(self):
        style = ttk.Style()
        try:
            style.theme_use("clam")
        except Exception:
            pass

        style.configure(
            "PW.Horizontal.TProgressbar",
            troughcolor="#E5E7EB",
            background=PURPLE2,
            bordercolor="#E5E7EB",
            lightcolor=PURPLE2,
            darkcolor=PURPLE2,
        )
        style.configure(
            "PW.TCombobox",
            fieldbackground=WHITE,
            background=WHITE,
            foreground=TEXT,
        )

    def clear(self, frame):
        for child in frame.winfo_children():
            child.destroy()

    def button(self, parent, text, command, color=PURPLE, width=None):
        b = tk.Button(
            parent,
            text=text,
            command=command,
            bg=color,
            fg=WHITE,
            activebackground=color,
            activeforeground=WHITE,
            relief="flat",
            bd=0,
            cursor="hand2",
            font=("Segoe UI", 10, "bold"),
            padx=16,
            pady=10,
        )
        if width:
            b.configure(width=width)
        b.bind("<Enter>", lambda e: b.configure(relief="raised", bd=1))
        b.bind("<Leave>", lambda e: b.configure(relief="flat", bd=0))
        return b

    def card(self, parent, color=WHITE):
        return tk.Frame(
            parent,
            bg=color,
            highlightbackground=BORDER,
            highlightthickness=1,
            bd=0,
        )

    def build_shell(self):
        # Header
        header = tk.Frame(self.root, bg=DARK, height=82)
        header.pack(fill="x")
        header.pack_propagate(False)

        logo = tk.Frame(header, bg=DARK)
        logo.pack(side="left", padx=25)

        tk.Label(
            logo, text="PDF", bg=DARK, fg=WHITE,
            font=("Segoe UI", 25, "bold")
        ).pack(side="left")
        tk.Label(
            logo, text=" WORLD", bg=DARK, fg="#A78BFA",
            font=("Segoe UI", 25, "bold")
        ).pack(side="left")
        tk.Label(
            logo, text="  by Sagnik", bg=DARK, fg="#CBD5E1",
            font=("Segoe UI", 10, "bold")
        ).pack(side="left", pady=(13, 0))

        self.header_status = tk.Label(
            header, textvariable=self.status,
            bg=DARK, fg="#CBD5E1",
            font=("Segoe UI", 9)
        )
        self.header_status.pack(side="right", padx=22)

        # Body
        body = tk.Frame(self.root, bg=BG)
        body.pack(fill="both", expand=True)

        # Sidebar
        self.sidebar = tk.Frame(body, bg=WHITE, width=215)
        self.sidebar.pack(side="left", fill="y")
        self.sidebar.pack_propagate(False)

        tk.Label(
            self.sidebar, text="TOOLS", bg=WHITE, fg=MUTED,
            font=("Segoe UI", 9, "bold")
        ).pack(anchor="w", padx=22, pady=(25, 8))

        nav = [
            ("⌂  Home", self.show_home),
            ("▣  Create PDF", self.show_create),
            ("↔  From PDF", self.show_from_pdf),
            ("✚  PDF Tools", self.show_pdf_tools),
        ]
        for text, command in nav:
            self.nav_button(text, command)

        tk.Frame(self.sidebar, bg=BORDER, height=1).pack(fill="x", padx=18, pady=18)

        self.nav_button("⚙  About", self.about)

        tk.Label(
            self.sidebar,
            text="100% local\nNo cloud uploads",
            bg=WHITE, fg=MUTED,
            justify="left",
            font=("Segoe UI", 8)
        ).pack(anchor="w", padx=22, pady=(25, 0))

        # Main area
        self.main = tk.Frame(body, bg=BG)
        self.main.pack(side="left", fill="both", expand=True)

        # Footer
        footer = tk.Frame(self.root, bg=DARK, height=36)
        footer.pack(fill="x")
        footer.pack_propagate(False)

        tk.Label(
            footer, text=f"{APP_NAME}  •  v{VERSION}",
            bg=DARK, fg="#A78BFA",
            font=("Segoe UI", 8, "bold")
        ).pack(side="left", padx=20)

        self.progress_bar = ttk.Progressbar(
            footer, style="PW.Horizontal.TProgressbar",
            variable=self.progress, maximum=100, length=210
        )
        self.progress_bar.pack(side="right", padx=(5, 20), pady=9)

    def nav_button(self, text, command):
        b = tk.Button(
            self.sidebar, text=text, command=command,
            bg=WHITE, fg=TEXT,
            activebackground="#F3E8FF", activeforeground=PURPLE,
            anchor="w", relief="flat", bd=0,
            cursor="hand2",
            font=("Segoe UI", 10, "bold"),
            padx=20, pady=11
        )
        b.pack(fill="x", padx=10, pady=2)

    # ----- Pages -----
    def page_title(self, title, subtitle):
        tk.Label(
            self.main, text=title, bg=BG, fg=TEXT,
            font=("Segoe UI", 25, "bold")
        ).pack(anchor="w", padx=30, pady=(28, 3))
        tk.Label(
            self.main, text=subtitle, bg=BG, fg=MUTED,
            font=("Segoe UI", 10)
        ).pack(anchor="w", padx=30, pady=(0, 20))

    def show_home(self):
        self.current_page.set("Home")
        self.clear(self.main)

        self.page_title(
            "Everything PDF, in one place.",
            "Fast local tools for creating, converting and managing PDF files."
        )

        hero = self.card(self.main, DARK)
        hero.pack(fill="x", padx=30, pady=(0, 20))

        left = tk.Frame(hero, bg=DARK)
        left.pack(side="left", fill="both", expand=True, padx=28, pady=24)

        tk.Label(
            left, text="PDF WORLD",
            bg=DARK, fg="#A78BFA",
            font=("Segoe UI", 11, "bold")
        ).pack(anchor="w")
        tk.Label(
            left, text="Your files stay on your computer.",
            bg=DARK, fg=WHITE,
            font=("Segoe UI", 19, "bold")
        ).pack(anchor="w", pady=(5, 4))
        tk.Label(
            left,
            text="Create PDFs, convert documents and use powerful PDF tools without uploading your files.",
            bg=DARK, fg="#CBD5E1",
            wraplength=600, justify="left",
            font=("Segoe UI", 10)
        ).pack(anchor="w")

        self.button(
            left, "＋  Create a PDF", self.show_create, PURPLE
        ).pack(anchor="w", pady=(18, 0))

        right = tk.Frame(hero, bg=DARK)
        right.pack(side="right", padx=28, pady=24)
        tk.Label(
            right, text="PDF", bg="#312E81", fg="#DDD6FE",
            font=("Segoe UI", 36, "bold"),
            padx=22, pady=15
        ).pack()

        grid = tk.Frame(self.main, bg=BG)
        grid.pack(fill="both", expand=True, padx=30)

        cards = [
            ("Create PDF", "Convert documents, images and text", self.show_create, PURPLE),
            ("From PDF", "Export to images, text and Office formats", self.show_from_pdf, BLUE),
            ("PDF Tools", "Merge, split, rotate, protect and more", self.show_pdf_tools, GREEN),
            ("Batch workflow", "Process several files into one output folder", self.show_create, ORANGE),
        ]

        for i, (title, desc, cmd, color) in enumerate(cards):
            r, c = divmod(i, 2)
            box = self.card(grid)
            box.grid(row=r, column=c, sticky="nsew", padx=7, pady=7)
            grid.grid_columnconfigure(c, weight=1)
            grid.grid_rowconfigure(r, weight=1)

            tk.Frame(box, bg=color, width=7).pack(side="left", fill="y")
            inner = tk.Frame(box, bg=WHITE)
            inner.pack(side="left", fill="both", expand=True, padx=18, pady=17)
            tk.Label(
                inner, text=title, bg=WHITE, fg=TEXT,
                font=("Segoe UI", 13, "bold")
            ).pack(anchor="w")
            tk.Label(
                inner, text=desc, bg=WHITE, fg=MUTED,
                font=("Segoe UI", 9),
                wraplength=380, justify="left"
            ).pack(anchor="w", pady=(5, 12))
            self.button(inner, "Open →", cmd, color).pack(anchor="w")

    def show_create(self):
        self.current_page.set("Create PDF")
        self.clear(self.main)
        self.page_title(
            "Create PDF",
            "Turn common documents, spreadsheets, presentations, images and text into PDF."
        )

        info = self.card(self.main)
        info.pack(fill="x", padx=30, pady=(0, 15))
        tk.Label(
            info,
            text="Tip: You can select multiple files. Each file becomes its own PDF.",
            bg=WHITE, fg=MUTED, font=("Segoe UI", 9)
        ).pack(anchor="w", padx=20, pady=12)

        grid = tk.Frame(self.main, bg=BG)
        grid.pack(fill="both", expand=True, padx=30)
        for c in range(2):
            grid.grid_columnconfigure(c, weight=1)

        actions = [
            ("📄", "Documents → PDF",
             "TXT • DOCX • HTML/HTM", self.convert_documents, PURPLE),
            ("📊", "Spreadsheets → PDF",
             "XLSX", self.convert_spreadsheets, GREEN),
            ("📽", "Presentations → PDF",
             "PPTX", self.convert_presentations, ORANGE),
            ("🖼", "Images → PDF",
             "PNG • JPG • JPEG • BMP • GIF • TIFF • WEBP",
             self.convert_images, BLUE),
            ("📚", "Images → One PDF",
             "Select several images and combine them in order.",
             self.images_to_one_pdf, PINK),
            ("↔", "Drop files here",
             "Drag & drop is available when tkinterdnd2 is installed.",
             self.show_drop_area, TEAL),
        ]

        for i, (icon, title, desc, cmd, color) in enumerate(actions):
            r, c = divmod(i, 2)
            box = self.card(grid)
            box.grid(row=r, column=c, sticky="nsew", padx=7, pady=7)
            grid.grid_rowconfigure(r, weight=1)

            tk.Label(
                box, text=icon, bg=WHITE, fg=color,
                font=("Segoe UI Emoji", 23)
            ).pack(anchor="w", padx=20, pady=(18, 3))
            tk.Label(
                box, text=title, bg=WHITE, fg=TEXT,
                font=("Segoe UI", 13, "bold")
            ).pack(anchor="w", padx=20)
            tk.Label(
                box, text=desc, bg=WHITE, fg=MUTED,
                font=("Segoe UI", 8.5),
                wraplength=380, justify="left"
            ).pack(anchor="w", padx=20, pady=(5, 12))
            self.button(box, "Choose files →", cmd, color).pack(
                anchor="w", padx=20, pady=(0, 18)
            )

    def show_from_pdf(self):
        self.current_page.set("From PDF")
        self.clear(self.main)
        self.page_title(
            "From PDF",
            "Extract or export PDF content. Visual exports preserve the original page appearance."
        )

        grid = tk.Frame(self.main, bg=BG)
        grid.pack(fill="both", expand=True, padx=30)

        actions = [
            ("🖼", "PDF → Images", "PNG pages at high resolution", self.pdf_to_images, BLUE),
            ("📝", "PDF → TXT", "Extract text only — no PDF filename is inserted", self.pdf_to_txt, CYAN),
            ("🌐", "PDF → HTML", "Create a clean HTML text export", self.pdf_to_html, PINK),
            ("📘", "PDF → Word", "DOCX with page images for layout fidelity", self.pdf_to_docx, PURPLE),
            ("📊", "PDF → Excel", "Extract text into worksheet rows", self.pdf_to_xlsx, GREEN),
            ("📽", "PDF → PowerPoint", "One original PDF page per slide", self.pdf_to_pptx, ORANGE),
        ]

        for i, (icon, title, desc, cmd, color) in enumerate(actions):
            r, c = divmod(i, 2)
            box = self.card(grid)
            box.grid(row=r, column=c, sticky="nsew", padx=7, pady=7)
            grid.grid_columnconfigure(c, weight=1)
            grid.grid_rowconfigure(r, weight=1)

            tk.Label(
                box, text=icon, bg=WHITE, fg=color,
                font=("Segoe UI Emoji", 23)
            ).pack(anchor="w", padx=20, pady=(18, 3))
            tk.Label(
                box, text=title, bg=WHITE, fg=TEXT,
                font=("Segoe UI", 13, "bold")
            ).pack(anchor="w", padx=20)
            tk.Label(
                box, text=desc, bg=WHITE, fg=MUTED,
                font=("Segoe UI", 8.5),
                wraplength=390, justify="left"
            ).pack(anchor="w", padx=20, pady=(5, 12))
            self.button(box, "Choose PDF →", cmd, color).pack(
                anchor="w", padx=20, pady=(0, 18)
            )

    def show_pdf_tools(self):
        self.current_page.set("PDF Tools")
        self.clear(self.main)
        self.page_title(
            "PDF Tools",
            "Organize, protect and modify PDF files locally."
        )

        grid = tk.Frame(self.main, bg=BG)
        grid.pack(fill="both", expand=True, padx=30)

        actions = [
            ("🔗", "Merge PDF", "Combine multiple PDFs in selected order", self.merge_pdfs, PURPLE),
            ("✂", "Split / Extract", "Create a PDF from selected page numbers", self.extract_pages, BLUE),
            ("↻", "Rotate Pages", "Rotate every page by 90°, 180° or 270°", self.rotate_pdf, CYAN),
            ("🗑", "Delete Pages", "Remove selected page numbers", self.delete_pages, RED),
            ("↕", "Reorder Pages", "Enter a new page order", self.reorder_pages, TEAL),
            ("🗜", "Compress PDF", "Compress page content streams where possible", self.compress_pdf, GREEN),
            ("💧", "Watermark", "Place a diagonal text watermark on every page", self.watermark_pdf, ORANGE),
            ("🔐", "Password Protect", "Encrypt a PDF with a password", self.protect_pdf, PINK),
            ("🔓", "Unlock PDF", "Remove encryption when you know the password", self.unlock_pdf, PURPLE),
        ]

        for i, (icon, title, desc, cmd, color) in enumerate(actions):
            r, c = divmod(i, 3)
            box = self.card(grid)
            box.grid(row=r, column=c, sticky="nsew", padx=6, pady=6)
            grid.grid_columnconfigure(c, weight=1)
            grid.grid_rowconfigure(r, weight=1)

            tk.Label(
                box, text=icon, bg=WHITE, fg=color,
                font=("Segoe UI Emoji", 20)
            ).pack(anchor="w", padx=16, pady=(15, 2))
            tk.Label(
                box, text=title, bg=WHITE, fg=TEXT,
                font=("Segoe UI", 11, "bold")
            ).pack(anchor="w", padx=16)
            tk.Label(
                box, text=desc, bg=WHITE, fg=MUTED,
                font=("Segoe UI", 8),
                wraplength=260, justify="left"
            ).pack(anchor="w", padx=16, pady=(4, 10))
            self.button(box, "Open →", cmd, color).pack(
                anchor="w", padx=16, pady=(0, 15)
            )

    def show_drop_area(self):
        self.clear(self.main)
        self.page_title(
            "Drop files",
            "Drag files into the box, then choose the output type."
        )

        box = self.card(self.main)
        box.pack(fill="both", expand=True, padx=30, pady=(0, 30))

        label = tk.Label(
            box,
            text="DROP FILES HERE\n\nor click to choose files",
            bg="#F5F3FF", fg=PURPLE,
            font=("Segoe UI", 18, "bold"),
            justify="center",
            relief="groove",
            bd=1,
            cursor="hand2"
        )
        label.pack(fill="both", expand=True, padx=30, pady=30)
        label.bind("<Button-1>", lambda e: self.handle_drop_click())

        if DND_OK:
            try:
                label.drop_target_register(DND_FILES)
                label.dnd_bind("<<Drop>>", self.handle_drop_event)
                label.configure(text="DROP FILES HERE\n\nDrag & drop is enabled")
            except Exception:
                pass
        else:
            label.configure(
                text="DROP FILES HERE\n\nClick to choose files\n\n"
                     "(Install tkinterdnd2 to enable real drag & drop.)"
            )

        self.button(
            self.main, "← Back", self.show_create, DARK
        ).pack(anchor="w", padx=30, pady=(0, 25))

    # ----- Conversion selection -----
    def pick_many(self, filetypes, title):
        return list(filedialog.askopenfilenames(title=title, filetypes=filetypes))

    def output_for_each(self, files):
        out = choose_output_dir()
        if not out:
            return None
        return Path(out)

    def run_batch(self, files, worker, success_text="Done"):
        files = [Path(f) for f in files]
        if not files:
            return

        def task():
            errors = []
            total = len(files)
            for i, src in enumerate(files, 1):
                try:
                    worker(src)
                except Exception as exc:
                    errors.append(f"{src.name}: {exc}")
                self.root.after(
                    0, lambda p=(i / total) * 100:
                    self.progress.set(p)
                )

            def finish():
                self.progress.set(100)
                if errors:
                    self.status.set(f"Finished with {len(errors)} error(s)")
                    detail = "\n".join(errors[:12])
                    if len(errors) > 12:
                        detail += f"\n…and {len(errors) - 12} more."
                    messagebox.showwarning(
                        "PDF World — some files failed",
                        detail
                    )
                else:
                    self.status.set(success_text)
                    messagebox.showinfo("PDF World", success_text)
                self.progress.set(0)

            self.root.after(0, finish)

        threading.Thread(target=task, daemon=True).start()

    def convert_documents(self):
        files = self.pick_many(
            [("Documents", "*.txt *.docx *.html *.htm"),
             ("All files", "*.*")],
            "Choose documents"
        )
        if not files:
            return
        out = self.output_for_each(files)
        if not out:
            return

        def worker(src):
            dest = unique_path(out / f"{src.stem}.pdf")
            ext = src.suffix.lower()
            if ext == ".txt":
                self.builder.txt_to_pdf(src, dest)
            elif ext == ".docx":
                self.builder.docx_to_pdf(src, dest)
            elif ext in (".html", ".htm"):
                self.builder.html_to_pdf(src, dest)
            else:
                raise ValueError("Unsupported document type")

        self.run_batch(files, worker, "Documents converted successfully.")

    def convert_spreadsheets(self):
        files = self.pick_many(
            [("Excel", "*.xlsx"), ("All files", "*.*")],
            "Choose spreadsheets"
        )
        if not files:
            return
        out = self.output_for_each(files)
        if not out:
            return
        self.run_batch(
            files,
            lambda src: self.builder.xlsx_to_pdf(
                src, unique_path(out / f"{src.stem}.pdf")
            ),
            "Spreadsheets converted successfully."
        )

    def convert_presentations(self):
        files = self.pick_many(
            [("PowerPoint", "*.pptx"), ("All files", "*.*")],
            "Choose presentations"
        )
        if not files:
            return
        out = self.output_for_each(files)
        if not out:
            return
        self.run_batch(
            files,
            lambda src: self.builder.pptx_to_pdf(
                src, unique_path(out / f"{src.stem}.pdf")
            ),
            "Presentations converted successfully."
        )

    def convert_images(self):
        files = self.pick_many(
            [("Images", "*.png *.jpg *.jpeg *.bmp *.gif *.tif *.tiff *.webp"),
             ("All files", "*.*")],
            "Choose images"
        )
        if not files:
            return
        out = self.output_for_each(files)
        if not out:
            return
        self.run_batch(
            files,
            lambda src: self.builder.image_to_pdf(
                src, unique_path(out / f"{src.stem}.pdf")
            ),
            "Images converted successfully."
        )

    def images_to_one_pdf(self):
        files = self.pick_many(
            [("Images", "*.png *.jpg *.jpeg *.bmp *.gif *.tif *.tiff *.webp")],
            "Choose images in desired order"
        )
        if not files:
            return
        dest = choose_save(".pdf", "PDF", "images.pdf")
        if not dest:
            return
        def task():
            try:
                image_paths = [Path(f) for f in files]
                total = len(image_paths)
                # Build the combined PDF in exactly the order selected by the user.
                self.builder.images_to_pdf(image_paths, Path(dest))
                self.root.after(0, lambda: self.progress.set(100))
                self.root.after(0, lambda: self.status.set("Combined PDF ready"))
                self.root.after(0, lambda: messagebox.showinfo(
                    "PDF World", "Combined PDF created successfully."
                ))
            except Exception as exc:
                self.root.after(0, lambda: messagebox.showerror(
                    "PDF World", str(exc)
                ))
            finally:
                self.root.after(0, lambda: self.progress.set(0))

        threading.Thread(target=task, daemon=True).start()

    # ----- From PDF -----
    def pdf_to_images(self):
        src = choose_pdf()
        if not src:
            return
        out = choose_output_dir()
        if not out:
            return

        def task():
            errors = []
            try:
                pdf = fitz.open(src)
                total = len(pdf)
                for i, page in enumerate(pdf):
                    pix = page.get_pixmap(
                        matrix=fitz.Matrix(2.0, 2.0),
                        alpha=False
                    )
                    target = unique_path(
                        Path(out) / f"{Path(src).stem}_page_{i + 1}.png"
                    )
                    pix.save(str(target))
                    self.root.after(
                        0, lambda p=((i + 1) / total) * 100:
                        self.progress.set(p)
                    )
                pdf.close()
            except Exception as exc:
                errors.append(str(exc))

            def finish():
                self.progress.set(0)
                if errors:
                    messagebox.showerror("PDF World", errors[0])
                else:
                    self.status.set("PDF exported as images")
                    messagebox.showinfo("PDF World", "Images created successfully.")
            self.root.after(0, finish)

        threading.Thread(target=task, daemon=True).start()

    def pdf_to_txt(self):
        src = choose_pdf()
        if not src:
            return
        dest = choose_save(".txt", "Text", Path(src).stem + ".txt")
        if not dest:
            return

        def task():
            try:
                pdf = fitz.open(src)
                # IMPORTANT: only extracted PDF text is written.
                # The PDF filename is never inserted into the TXT.
                text = "\n\n".join(page.get_text("text") for page in pdf)
                pdf.close()
                Path(dest).write_text(text, encoding="utf-8")
                self.root.after(0, lambda: (
                    self.status.set("Text extracted successfully"),
                    messagebox.showinfo("PDF World", "TXT created successfully.")
                ))
            except Exception as exc:
                self.root.after(0, lambda: messagebox.showerror(
                    "PDF World", str(exc)
                ))

        threading.Thread(target=task, daemon=True).start()

    def pdf_to_html(self):
        src = choose_pdf()
        if not src:
            return
        dest = choose_save(".html", "HTML", Path(src).stem + ".html")
        if not dest:
            return

        def task():
            try:
                pdf = fitz.open(src)
                parts = [
                    "<!doctype html><html><head><meta charset='utf-8'>",
                    "<title>PDF export</title>",
                    "<style>body{font-family:Segoe UI,Arial,sans-serif;"
                    "max-width:900px;margin:40px auto;line-height:1.6}"
                    ".page{margin-bottom:40px;padding-bottom:20px;"
                    "border-bottom:1px solid #ddd}</style></head><body>"
                ]
                for i, page in enumerate(pdf):
                    parts.append(f"<section class='page'><h2>Page {i+1}</h2>")
                    parts.append(
                        "<pre style='white-space:pre-wrap'>"
                        + html.escape(page.get_text("text"))
                        + "</pre></section>"
                    )
                pdf.close()
                parts.append("</body></html>")
                Path(dest).write_text("".join(parts), encoding="utf-8")
                self.root.after(0, lambda: (
                    self.status.set("HTML created successfully"),
                    messagebox.showinfo("PDF World", "HTML created successfully.")
                ))
            except Exception as exc:
                self.root.after(0, lambda: messagebox.showerror(
                    "PDF World", str(exc)
                ))

        threading.Thread(target=task, daemon=True).start()

    def pdf_to_docx(self):
        src = choose_pdf()
        if not src:
            return
        dest = choose_save(".docx", "Word document", Path(src).stem + ".docx")
        if not dest:
            return

        def task():
            temp_dir = Path(dest).with_name("__pw_docx_images")
            temp_dir.mkdir(exist_ok=True)
            try:
                pdf = fitz.open(src)
                doc = Document()
                for i, page in enumerate(pdf):
                    pix = page.get_pixmap(
                        matrix=fitz.Matrix(1.5, 1.5),
                        alpha=False
                    )
                    image_path = temp_dir / f"page_{i+1}.png"
                    pix.save(str(image_path))
                    # Preserve the visual page exactly as an image.
                    section = doc.sections[0] if i == 0 else doc.add_section()
                    section.top_margin = Inches(0.35)
                    section.bottom_margin = Inches(0.35)
                    section.left_margin = Inches(0.35)
                    section.right_margin = Inches(0.35)
                    doc.add_picture(str(image_path), width=Inches(7.25))
                    if i != len(pdf) - 1:
                        doc.add_page_break()
                    self.root.after(
                        0, lambda p=((i + 1) / len(pdf)) * 100:
                        self.progress.set(p)
                    )
                pdf.close()
                doc.save(dest)
            finally:
                for p in temp_dir.glob("*"):
                    try:
                        p.unlink()
                    except OSError:
                        pass
                try:
                    temp_dir.rmdir()
                except OSError:
                    pass

            self.root.after(0, lambda: (
                self.progress.set(0),
                self.status.set("Word export complete"),
                messagebox.showinfo(
                    "PDF World",
                    "Word document created.\n\n"
                    "Pages are preserved as images for maximum visual fidelity."
                )
            ))

        threading.Thread(target=task, daemon=True).start()

    def pdf_to_xlsx(self):
        src = choose_pdf()
        if not src:
            return
        dest = choose_save(".xlsx", "Excel workbook", Path(src).stem + ".xlsx")
        if not dest:
            return

        def task():
            try:
                pdf = fitz.open(src)
                wb = Workbook()
                ws = wb.active
                ws.title = "PDF Text"
                row = 1

                for page_no, page in enumerate(pdf, 1):
                    ws.cell(row, 1, f"Page {page_no}")
                    ws.cell(row, 1).font = Font(bold=True)
                    row += 1
                    for line in page.get_text("text").splitlines():
                        ws.cell(row, 1, line)
                        row += 1
                    row += 1

                ws.column_dimensions["A"].width = 100
                wb.save(dest)
                pdf.close()

                self.root.after(0, lambda: (
                    self.status.set("Excel export complete"),
                    messagebox.showinfo(
                        "PDF World",
                        "Excel workbook created.\n\n"
                        "Text is exported into worksheet rows."
                    )
                ))
            except Exception as exc:
                self.root.after(0, lambda: messagebox.showerror(
                    "PDF World", str(exc)
                ))

        threading.Thread(target=task, daemon=True).start()

    def pdf_to_pptx(self):
        src = choose_pdf()
        if not src:
            return
        dest = choose_save(".pptx", "PowerPoint", Path(src).stem + ".pptx")
        if not dest:
            return

        def task():
            temp_dir = Path(dest).with_name("__pw_ppt_images")
            temp_dir.mkdir(exist_ok=True)
            try:
                pdf = fitz.open(src)
                prs = Presentation()
                prs.slide_width = PptInches(13.333)
                prs.slide_height = PptInches(7.5)

                # Remove default slide layouts by using the blank layout.
                blank = prs.slide_layouts[6]

                for i, page in enumerate(pdf):
                    pix = page.get_pixmap(
                        matrix=fitz.Matrix(1.5, 1.5),
                        alpha=False
                    )
                    img = temp_dir / f"page_{i+1}.png"
                    pix.save(str(img))

                    slide = prs.slides.add_slide(blank)
                    slide.shapes.add_picture(
                        str(img), 0, 0,
                        width=prs.slide_width,
                        height=prs.slide_height
                    )

                    self.root.after(
                        0, lambda p=((i + 1) / len(pdf)) * 100:
                        self.progress.set(p)
                    )

                pdf.close()
                prs.save(dest)
            finally:
                for p in temp_dir.glob("*"):
                    try:
                        p.unlink()
                    except OSError:
                        pass
                try:
                    temp_dir.rmdir()
                except OSError:
                    pass

            self.root.after(0, lambda: (
                self.progress.set(0),
                self.status.set("PowerPoint export complete"),
                messagebox.showinfo(
                    "PDF World",
                    "PowerPoint created.\n\n"
                    "Each PDF page is preserved as a full-slide image."
                )
            ))

        threading.Thread(target=task, daemon=True).start()

    # ----- PDF tools -----
    def merge_pdfs(self):
        files = self.pick_many([("PDF", "*.pdf")], "Choose PDFs in merge order")
        if len(files) < 2:
            if files:
                messagebox.showwarning("PDF World", "Choose at least 2 PDFs.")
            return
        dest = choose_save(".pdf", "PDF", "merged.pdf")
        if not dest:
            return

        def task():
            try:
                PDFTools.merge([Path(f) for f in files], Path(dest))
                self.root.after(0, lambda: messagebox.showinfo(
                    "PDF World", "PDFs merged successfully."
                ))
            except Exception as exc:
                self.root.after(0, lambda: messagebox.showerror(
                    "PDF World", str(exc)
                ))
        threading.Thread(target=task, daemon=True).start()

    def extract_pages(self):
        src = choose_pdf()
        if not src:
            return
        reader = PdfReader(src)
        total = len(reader.pages)
        value = simpledialog.askstring(
            "Extract pages",
            f"PDF has {total} pages.\nEnter pages/ranges, e.g.:\n1,3,5-8"
        )
        if not value:
            return
        pages = self.parse_page_list(value, total)
        if not pages:
            messagebox.showerror("PDF World", "No valid pages were entered.")
            return
        dest = choose_save(".pdf", "PDF", Path(src).stem + "_extracted.pdf")
        if not dest:
            return
        try:
            PDFTools.extract(Path(src), Path(dest), pages)
            messagebox.showinfo("PDF World", "Pages extracted successfully.")
        except Exception as exc:
            messagebox.showerror("PDF World", str(exc))

    def parse_page_list(self, value, total):
        result = []
        for part in value.replace(" ", "").split(","):
            if not part:
                continue
            if "-" in part:
                try:
                    a, b = part.split("-", 1)
                    a, b = int(a), int(b)
                    if a > b:
                        a, b = b, a
                    result.extend(range(max(1, a), min(total, b) + 1))
                except ValueError:
                    continue
            else:
                try:
                    n = int(part)
                    if 1 <= n <= total:
                        result.append(n)
                except ValueError:
                    continue
        # Keep order but remove duplicates.
        return list(dict.fromkeys(result))

    def rotate_pdf(self):
        src = choose_pdf()
        if not src:
            return
        degrees = simpledialog.askinteger(
            "Rotate PDF",
            "Rotation: enter 90, 180 or 270",
            minvalue=90, maxvalue=270
        )
        if degrees not in (90, 180, 270):
            return
        dest = choose_save(".pdf", "PDF", Path(src).stem + f"_rotated_{degrees}.pdf")
        if not dest:
            return
        try:
            PDFTools.rotate(Path(src), Path(dest), degrees)
            messagebox.showinfo("PDF World", "PDF rotated successfully.")
        except Exception as exc:
            messagebox.showerror("PDF World", str(exc))

    def delete_pages(self):
        src = choose_pdf()
        if not src:
            return
        reader = PdfReader(src)
        total = len(reader.pages)
        value = simpledialog.askstring(
            "Delete pages",
            f"PDF has {total} pages.\nEnter pages/ranges to delete, e.g. 2,5-7"
        )
        if not value:
            return
        delete_set = set(self.parse_page_list(value, total))
        if not delete_set:
            messagebox.showerror("PDF World", "No valid pages were entered.")
            return
        if len(delete_set) >= total:
            messagebox.showerror("PDF World", "You cannot delete every page.")
            return
        dest = choose_save(".pdf", "PDF", Path(src).stem + "_pages_removed.pdf")
        if not dest:
            return
        try:
            PDFTools.delete_pages(Path(src), Path(dest), delete_set)
            messagebox.showinfo("PDF World", "Pages deleted successfully.")
        except Exception as exc:
            messagebox.showerror("PDF World", str(exc))

    def reorder_pages(self):
        src = choose_pdf()
        if not src:
            return
        total = len(PdfReader(src).pages)
        value = simpledialog.askstring(
            "Reorder pages",
            f"PDF has {total} pages.\nEnter complete order, e.g. 3,1,2,4"
        )
        if not value:
            return
        order = self.parse_page_list(value, total)
        if len(order) != total or set(order) != set(range(1, total + 1)):
            messagebox.showerror(
                "PDF World",
                "Enter every page exactly once."
            )
            return
        dest = choose_save(".pdf", "PDF", Path(src).stem + "_reordered.pdf")
        if not dest:
            return
        try:
            PDFTools.reorder(Path(src), Path(dest), order)
            messagebox.showinfo("PDF World", "Pages reordered successfully.")
        except Exception as exc:
            messagebox.showerror("PDF World", str(exc))

    def compress_pdf(self):
        src = choose_pdf()
        if not src:
            return
        dest = choose_save(".pdf", "PDF", Path(src).stem + "_compressed.pdf")
        if not dest:
            return
        try:
            PDFTools.compress(Path(src), Path(dest))
            before = Path(src).stat().st_size
            after = Path(dest).stat().st_size
            messagebox.showinfo(
                "PDF World",
                f"Compression finished.\n\n"
                f"Original: {before / 1024:.1f} KB\n"
                f"Output:   {after / 1024:.1f} KB"
            )
        except Exception as exc:
            messagebox.showerror("PDF World", str(exc))

    def watermark_pdf(self):
        src = choose_pdf()
        if not src:
            return
        text = simpledialog.askstring(
            "Watermark",
            "Watermark text:"
        )
        if not text:
            return
        dest = choose_save(".pdf", "PDF", Path(src).stem + "_watermarked.pdf")
        if not dest:
            return
        try:
            PDFTools.watermark(Path(src), Path(dest), text)
            messagebox.showinfo("PDF World", "Watermark added successfully.")
        except Exception as exc:
            messagebox.showerror("PDF World", str(exc))

    def protect_pdf(self):
        src = choose_pdf()
        if not src:
            return
        password = simpledialog.askstring(
            "Password protect",
            "Enter a password:",
            show="*"
        )
        if not password:
            return
        confirm = simpledialog.askstring(
            "Password protect",
            "Confirm password:",
            show="*"
        )
        if password != confirm:
            messagebox.showerror("PDF World", "Passwords do not match.")
            return
        dest = choose_save(".pdf", "PDF", Path(src).stem + "_protected.pdf")
        if not dest:
            return
        try:
            PDFTools.protect(Path(src), Path(dest), password)
            messagebox.showinfo("PDF World", "PDF protected successfully.")
        except Exception as exc:
            messagebox.showerror("PDF World", str(exc))

    def unlock_pdf(self):
        src = choose_pdf()
        if not src:
            return
        password = simpledialog.askstring(
            "Unlock PDF",
            "Enter the existing PDF password:",
            show="*"
        )
        if password is None:
            return
        dest = choose_save(".pdf", "PDF", Path(src).stem + "_unlocked.pdf")
        if not dest:
            return
        try:
            PDFTools.unlock(Path(src), Path(dest), password)
            messagebox.showinfo("PDF World", "PDF unlocked successfully.")
        except Exception as exc:
            messagebox.showerror("PDF World", str(exc))

    # ----- Drag/drop -----
    def handle_drop_click(self):
        files = self.pick_many(
            [("Supported", "*.txt *.docx *.xlsx *.pptx *.html *.htm "
                           "*.png *.jpg *.jpeg *.bmp *.gif *.tif *.tiff *.webp")],
            "Choose files"
        )
        if files:
            self.process_drop_files(files)

    def handle_drop_event(self, event):
        try:
            files = list(self.root.tk.splitlist(event.data))
        except Exception:
            files = []
        if files:
            self.process_drop_files(files)

    def process_drop_files(self, files):
        out = choose_output_dir()
        if not out:
            return
        out = Path(out)

        def worker(src):
            src = Path(src)
            ext = src.suffix.lower()
            dest = unique_path(out / f"{src.stem}.pdf")
            if ext == ".txt":
                self.builder.txt_to_pdf(src, dest)
            elif ext == ".docx":
                self.builder.docx_to_pdf(src, dest)
            elif ext == ".xlsx":
                self.builder.xlsx_to_pdf(src, dest)
            elif ext == ".pptx":
                self.builder.pptx_to_pdf(src, dest)
            elif ext in (".html", ".htm"):
                self.builder.html_to_pdf(src, dest)
            elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".gif",
                         ".tif", ".tiff", ".webp"):
                self.builder.image_to_pdf(src, dest)
            else:
                raise ValueError(f"Unsupported format: {ext}")

        self.run_batch(files, worker, "Dropped files converted successfully.")

    def about(self):
        messagebox.showinfo(
            f"{APP_NAME} — About",
            f"{APP_NAME}\nVersion {VERSION}\n\n"
            "A colourful local PDF toolkit built with Python.\n\n"
            "Designed to avoid the common TXT conversion problem where "
            "the input filename gets printed at the top of the PDF: "
            "PDF World v2 writes only the actual text content.\n\n"
            "Built-in conversion does not require LibreOffice.\n\n"
            "Important format note:\n"
            ".DOC/.XLS/.PPT are legacy binary Office formats and cannot "
            "be faithfully converted by these pure-Python converters. "
            "DOCX/XLSX/PPTX are supported. Complex Office layouts, "
            "macros and animations may not reproduce exactly.\n\n"
            "PDF → Word/PowerPoint preserve visual appearance using "
            "page images; PDF → Excel exports extracted text."
        )


def main():
    if DND_OK:
        root = TkinterDnD.Tk()
    else:
        root = tk.Tk()

    app = PDFWorldApp(root)

    # Global exception handler: show a friendly dialog instead of silently crashing.
    def report_exception(exc_type, exc_value, exc_tb):
        text = "".join(traceback.format_exception(exc_type, exc_value, exc_tb))
        print(text, file=sys.stderr)
        try:
            messagebox.showerror(
                "PDF World — unexpected error",
                f"{exc_value}\n\nA detailed error was written to the terminal."
            )
        except Exception:
            pass

    root.report_callback_exception = report_exception
    root.mainloop()


if __name__ == "__main__":
    main()
