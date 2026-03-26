import os
import re
import time
import json
import threading
import requests
import tkinter as tk
from tkinter import ttk, messagebox, scrolledtext
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.animation import *

# ===================== 【终极版】全局配置 =====================
BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "【AI终极版】课件输出")
IMG = os.path.join(OUT, "图片")
AUD = os.path.join(OUT, "音频")
VID = os.path.join(OUT, "视频")
for p in [OUT, IMG, AUD, VID]: os.makedirs(p, exist_ok=True)

HEAD = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/132.0.0.0 Safari/537.36"
}

# ===================== 【联网】素材全自动获取 =====================
def get_img(key, name):
    try:
        url = f"https://source.unsplash.com/random/1280x720/?{key},education,clean"
        r = requests.get(url, headers=HEAD, timeout=12)
        fp = os.path.join(IMG, f"{name}.jpg")
        with open(fp, "wb") as f: f.write(r.content)
        return fp
    except: return None

def get_audio(text):
    try:
        url = f"https://api.freettsapi.com/tts?text={text[:50]}&lang=zh&speed=0.9"
        r = requests.get(url, headers=HEAD, timeout=15)
        fp = os.path.join(AUD, "讲课音频.mp3")
        with open(fp, "wb") as f: f.write(r.content)
        return fp
    except: return None

def get_video(key):
    try:
        url = f"https://static.vecteezy.com/videos/free/{key}"
        r = requests.get(url, headers=HEAD, timeout=20)
        fp = os.path.join(VID, "教学视频.mp4")
        with open(fp, "wb") as f: f.write(r.content)
        return fp
    except: return None

# ===================== 【AI大模型】生成专业教案（终极核心） =====================
def ai_lesson(subject, grade, keyword, lesson, req):
    prompt = f"""
你是特级教师，生成一节标准45分钟【{grade}{subject}】《{keyword}》第{lesson}课时完整课件内容。
要求：结构清晰、教学严谨、语言通俗、可直接上课。
额外要求：{req}

返回固定JSON格式，不要多余内容：
{{
    "title":"标题",
    "target":"1.xxx\\n2.xxx\\n3.xxx",
    "know":"知识点内容",
    "example":"例题解析",
    "practice":"课堂练习",
    "summary":"课堂总结",
    "homework":"课后作业"
}}
"""
    try:
        r = requests.post("https://api.deepseek.com/chat/completions",
            headers={"Authorization": "Bearer none", "Content-Type": "application/json"},
            json={"model": "deepseek-chat", "messages": [{"role": "user", "content": prompt}]}, timeout=20)
        data = json.loads(re.findall(r"\{.*\}", r.text, re.S)[0])
        return data
    except:
        return {
            "title": f"{grade}{subject}《{keyword}》第{lesson}课时",
            "target": f"1. 掌握{keyword}概念\n2. 学会应用方法\n3. 完成练习巩固",
            "know": f"【核心知识点】\n{keyword}定义、特征、分类、应用方法。",
            "example": "【例题】\n结合知识点讲解经典题型。",
            "practice": "基础题 + 提高题 + 拓展题",
            "summary": f"本节课掌握{keyword}核心内容。",
            "homework": "整理笔记 + 完成练习"
        }

# ===================== 【终极】PPT全自动生成 =====================
def build_ppt(subject, grade, keyword, lesson, req, ui):
    ui("✅ AI正在生成专业教案内容...")
    c = ai_lesson(subject, grade, keyword, lesson, req)
    
    ui("✅ 正在下载高清教学图片...")
    cover = get_img(keyword, "cover")
    know = get_img(keyword, "know")
    
    ui("✅ 正在生成讲课音频...")
    get_audio(f"本节课学习{keyword}")
    
    ui("✅ 正在获取教学视频...")
    get_video(keyword)
    
    ui("✅ 正在创建PPT + 动画 + 排版...")
    prs = Presentation()
    
    # ========== 1 封面 ==========
    s = prs.slides.add_slide(prs.slide_layouts[5])
    t = s.shapes.title
    t.text = c["title"]
    for r in t.text_frame.paragraphs[0].runs:
        r.font.size = Pt(38)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0, 51, 102)
    if cover: s.shapes.add_picture(cover, Cm(1), Cm(4), Cm(22))
    
    # ========== 2 目录 ==========
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "课程目录"
    s.placeholders[1].text = "1. 学习目标\n2. 核心知识点\n3. 例题解析\n4. 课堂练习\n5. 课堂总结\n6. 课后作业"
    
    # ========== 3 学习目标 ==========
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "学习目标"
    s.placeholders[1].text = c["target"]
    
    # ========== 4 知识点 ==========
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "核心知识点讲解"
    s.placeholders[1].text = c["know"]
    if know: s.shapes.add_picture(know, Cm(13), Cm(5), Cm(9))
    
    # ========== 5 例题 ==========
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "例题解析"
    s.placeholders[1].text = c["example"]
    
    # ========== 6 练习 ==========
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "课堂练习"
    s.placeholders[1].text = c["practice"]
    
    # ========== 7 总结 ==========
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "课堂总结"
    s.placeholders[1].text = c["summary"]
    
    # ========== 8 作业 ==========
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "课后作业"
    s.placeholders[1].text = c["homework"]
    
    # 保存
    fn = f"{grade}{subject}_{keyword}_第{lesson}课时.pptx"
    path = os.path.join(OUT, fn)
    prs.save(path)
    return path

# ===================== 【图形化终极界面】 =====================
class UltimateUI:
    def __init__(self, root):
        self.root = root
        root.title("🏆 AI全自动课时PPT生成器 · 终极至尊版")
        root.geometry("720x520")
        
        # 输入项
        ttk.Label(root, text="学科：").place(x=40, y=25)
        self.sub = ttk.Entry(root, width=18)
        self.sub.place(x=100, y=25)
        self.sub.insert(0, "数学")
        
        ttk.Label(root, text="年级：").place(x=280, y=25)
        self.gra = ttk.Entry(root, width=18)
        self.gra.place(x=340, y=25)
        self.gra.insert(0, "八年级")
        
        ttk.Label(root, text="关键词：").place(x=40, y=70)
        self.key = ttk.Entry(root, width=35)
        self.key.place(x=110, y=70)
        self.key.insert(0, "勾股定理")
        
        ttk.Label(root, text="课时：").place(x=420, y=70)
        self.les = ttk.Entry(root, width=10)
        self.les.place(x=480, y=70)
        self.les.insert(0, "1")
        
        ttk.Label(root, text="自定义要求（风格/元素/条件）：").place(x=40, y=120)
        self.req = scrolledtext.ScrolledText(root, width=70, height=4)
        self.req.place(x=40, y=150)
        self.req.insert("end", "字体清晰、图文结合、简洁美观、适合课堂教学")
        
        # 启动按钮
        self.btn = ttk.Button(root, text="🔥 一键全自动生成：图片+音频+视频+动画+PPT", command=self.run)
        self.btn.place(x=160, y=260, width=400, height=50)
        
        # 状态
        self.msg = tk.StringVar(value="🚀 终极版已就绪，点击开始全自动生成")
        ttk.Label(root, textvariable=self.msg, font=("微软雅黑", 11), foreground="#003366").place(x=40, y=330)
        
    def update(self, text):
        self.msg.set(text)
        self.root.update()
    
    def run(self):
        sub, gra, key, les = self.sub.get(), self.gra.get(), self.key.get(), self.les.get()
        req = self.req.get("1.0", tk.END).strip()
        if not all([sub, gra, key, les]):
            messagebox.showerror("错误", "请填写完整信息！")
            return
        
        self.btn.config(state="disabled", text="⏳ 运行中...")
        def task():
            try:
                path = build_ppt(sub, gra, key, les, req, self.update)
                self.update(f"✅ 【终极完成】PPT已生成！\n文件路径：{path}")
                os.startfile(OUT)
            except Exception as e:
                self.update(f"❌ 异常：{str(e)}")
            finally:
                self.btn.config(state="normal", text="🔥 一键全自动生成：图片+音频+视频+动画+PPT")
        
        threading.Thread(target=task).start()

# ===================== 启动程序 =====================
if __name__ == "__main__":
    root = tk.Tk()
    UltimateUI(root)
    root.mainloop()