import sys
import os
from datetime import datetime
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QTabWidget, QTextEdit, QTableWidget, QTableWidgetItem, QPushButton,
    QFileDialog, QLabel, QSlider, QMessageBox, QSplitter, QToolBar,
    QAction, QStatusBar, QComboBox, QSpinBox, QLineEdit,
    QFrame, QGridLayout, QGroupBox, QProgressBar, QGraphicsDropShadowEffect,
    QDialog, QDesktopWidget, QToolButton
)
from PyQt5.QtCore import Qt, QUrl, QTimer, pyqtSignal, QSize, QPropertyAnimation, QEasingCurve, QRect
from PyQt5.QtGui import QFont, QIcon, QColor, QPalette, QLinearGradient, QBrush, QPixmap, QFontDatabase, QPainter, QPen
from PyQt5.QtMultimedia import QMediaPlayer, QMediaContent
from PyQt5.QtMultimediaWidgets import QVideoWidget

# 文档处理库
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import Workbook, load_workbook
import pandas as pd

# 视频编辑 (moviepy 1.0.3)
from moviepy.editor import VideoFileClip

# =============================== 全局样式（极致美观版） ===============================
STYLE_SHEET = """
/* 主窗口：动态渐变背景 */
QMainWindow {
    background: qlineargradient(x1: 0, y1: 0, x2: 1, y2: 1,
                                stop: 0 #f5f7fa, stop: 0.5 #eef2f6, stop: 1 #e8ecf2);
}

/* 选项卡面板：毛玻璃效果 */
QTabWidget::pane {
    border: none;
    border-radius: 28px;
    background: rgba(255, 255, 255, 0.7);
    margin: 8px;
    padding: 12px;
}

QTabBar::tab {
    background: rgba(250, 250, 252, 0.6);
    border-top-left-radius: 24px;
    border-top-right-radius: 24px;
    padding: 12px 32px;
    margin-right: 8px;
    font-weight: 600;
    font-size: 13px;
    color: #2c3e50;
    transition: 0.3s cubic-bezier(0.4, 0, 0.2, 1);
}
QTabBar::tab:selected {
    background: white;
    border-bottom: 3px solid #6366f1;
    color: #6366f1;
}
QTabBar::tab:hover {
    background: rgba(255, 255, 255, 0.9);
    transform: translateY(-2px);
}

/* 按钮：精致圆角 + 阴影 + 动画 */
QPushButton, QToolButton {
    background: white;
    border: none;
    border-radius: 24px;
    padding: 8px 20px;
    font-weight: 600;
    color: #2c3e50;
    min-width: 80px;
    font-size: 12px;
}
QPushButton:hover, QToolButton:hover {
    background: #f8fafc;
    transform: translateY(-2px);
}
QPushButton:pressed, QToolButton:pressed {
    background: #eef2f6;
    transform: translateY(0px);
}

/* 工具栏按钮样式 */
QToolBar QToolButton {
    background: transparent;
    border: none;
    padding: 8px 16px;
    border-radius: 24px;
}
QToolBar QToolButton:hover {
    background: rgba(0, 0, 0, 0.05);
}

/* 输入框、表格、文本编辑器 */
QTextEdit, QTableWidget, QLineEdit, QSpinBox, QComboBox {
    background: rgba(255, 255, 255, 0.95);
    border: 1px solid #e2e8f0;
    border-radius: 20px;
    padding: 10px 14px;
    selection-background-color: #6366f1;
    selection-color: white;
    font-size: 12px;
}
QTextEdit:focus, QTableWidget:focus, QLineEdit:focus, QSpinBox:focus, QComboBox:focus {
    border-color: #6366f1;
    outline: none;
    background: white;
}

/* 滚动条美化 */
QScrollBar:vertical {
    background: #eef2f6;
    width: 8px;
    border-radius: 4px;
}
QScrollBar::handle:vertical {
    background: #cbd5e1;
    border-radius: 4px;
    min-height: 30px;
}
QScrollBar::handle:vertical:hover {
    background: #94a3b8;
}

/* 分组框 */
QGroupBox {
    border: 1px solid #e2e8f0;
    border-radius: 28px;
    margin-top: 20px;
    padding-top: 20px;
    background: rgba(255, 255, 255, 0.5);
}
QGroupBox::title {
    subcontrol-origin: margin;
    left: 20px;
    padding: 0 12px;
    color: #2c3e50;
    font-weight: 700;
    font-size: 13px;
}

/* 状态栏 */
QStatusBar {
    background: rgba(255, 255, 255, 0.6);
    border-top: 1px solid rgba(226, 232, 240, 0.8);
    border-radius: 0 0 28px 28px;
    padding: 4px 16px;
    color: #475569;
}

/* 标签 */
QLabel {
    color: #2c3e50;
    font-weight: 500;
    font-size: 12px;
}

/* 进度条 */
QProgressBar {
    border: none;
    background: #eef2f6;
    border-radius: 12px;
    text-align: center;
    height: 6px;
}
QProgressBar::chunk {
    background: #6366f1;
    border-radius: 12px;
}

/* 菜单栏 */
QMenuBar {
    background: transparent;
    color: #2c3e50;
    padding: 4px;
}
QMenuBar::item {
    padding: 6px 12px;
    border-radius: 12px;
}
QMenuBar::item:selected {
    background: rgba(99, 102, 241, 0.1);
}
QMenu {
    background: white;
    border: 1px solid #e2e8f0;
    border-radius: 16px;
    padding: 8px;
}
QMenu::item {
    padding: 8px 24px;
    border-radius: 12px;
}
QMenu::item:selected {
    background: #f1f5f9;
}
"""

def add_shadow(widget, blur_radius=20, offset=(0, 4), color=Qt.gray, alpha=60):
    """为控件添加阴影效果"""
    shadow = QGraphicsDropShadowEffect()
    shadow.setBlurRadius(blur_radius)
    shadow.setOffset(offset[0], offset[1])
    shadow.setColor(QColor(0, 0, 0, alpha))
    widget.setGraphicsEffect(shadow)

# =============================== PPT 全屏放映窗口 ===============================
class SlideShowWindow(QDialog):
    """全屏幻灯片放映窗口"""
    def __init__(self, presentation, parent=None):
        super().__init__(parent)
        self.presentation = presentation
        self.current_slide_index = 0
        self.total_slides = len(presentation.slides)
        self.init_ui()

    def init_ui(self):
        self.setWindowFlags(Qt.FramelessWindowHint | Qt.WindowStaysOnTopHint)
        self.setWindowState(Qt.WindowFullScreen)
        self.setStyleSheet("""
            background: qlineargradient(x1: 0, y1: 0, x2: 1, y2: 1,
                                        stop: 0 #1e293b, stop: 1 #0f172a);
        """)

        # 主布局
        layout = QVBoxLayout()
        layout.setContentsMargins(60, 60, 60, 60)

        # 内容标签
        self.content_label = QLabel()
        self.content_label.setAlignment(Qt.AlignCenter)
        self.content_label.setStyleSheet("""
            color: white; 
            background: rgba(255,255,255,0.1);
            border-radius: 32px;
            padding: 60px;
            font-size: 32px;
            font-weight: 500;
        """)
        self.content_label.setWordWrap(True)
        layout.addWidget(self.content_label)

        # 页码标签
        self.page_label = QLabel()
        self.page_label.setAlignment(Qt.AlignRight | Qt.AlignBottom)
        self.page_label.setStyleSheet("""
            color: rgba(255,255,255,0.5); 
            background: transparent; 
            font-size: 14px; 
            padding: 20px;
            font-weight: 500;
        """)
        layout.addWidget(self.page_label)

        self.setLayout(layout)
        self.update_slide()

    def update_slide(self):
        """更新当前幻灯片内容"""
        slide = self.presentation.slides[self.current_slide_index]
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        content = "\n\n".join(texts) if texts else "（此幻灯片无文本内容）"
        self.content_label.setText(content)
        self.page_label.setText(f"{self.current_slide_index+1} / {self.total_slides}")

    def mousePressEvent(self, event):
        """鼠标点击：左键下一页，右键上一页"""
        if event.button() == Qt.LeftButton:
            self.next_slide()
        elif event.button() == Qt.RightButton:
            self.prev_slide()

    def keyPressEvent(self, event):
        """键盘控制"""
        if event.key() == Qt.Key_Right:
            self.next_slide()
        elif event.key() == Qt.Key_Left:
            self.prev_slide()
        elif event.key() == Qt.Key_Escape:
            self.close()

    def next_slide(self):
        if self.current_slide_index < self.total_slides - 1:
            self.current_slide_index += 1
            self.update_slide()

    def prev_slide(self):
        if self.current_slide_index > 0:
            self.current_slide_index -= 1
            self.update_slide()

# =============================== 文字编辑器（修复UI bug） ===============================
class WordEditor(QWidget):
    def __init__(self):
        super().__init__()
        self.current_file = None
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout()
        layout.setContentsMargins(24, 24, 24, 24)
        layout.setSpacing(16)
        
        # 工具栏 - 使用 QToolBar 和 QAction 确保文字显示
        toolbar = QToolBar()
        toolbar.setMovable(False)
        toolbar.setIconSize(QSize(20, 20))
        toolbar.setStyleSheet("QToolBar { spacing: 4px; }")
        
        # 使用 QAction 并设置文本，确保文字显示
        new_action = QAction("📄 新建", self)
        new_action.triggered.connect(self.new_file)
        toolbar.addAction(new_action)
        
        open_action = QAction("📂 打开", self)
        open_action.triggered.connect(self.open_file)
        toolbar.addAction(open_action)
        
        save_action = QAction("💾 保存", self)
        save_action.triggered.connect(self.save_file)
        toolbar.addAction(save_action)
        
        saveas_action = QAction("📎 另存为", self)
        saveas_action.triggered.connect(self.save_as_file)
        toolbar.addAction(saveas_action)
        
        toolbar.addSeparator()
        
        bold_action = QAction("🔲 加粗", self)
        bold_action.triggered.connect(self.bold)
        toolbar.addAction(bold_action)
        
        italic_action = QAction("📐 斜体", self)
        italic_action.triggered.connect(self.italic)
        toolbar.addAction(italic_action)
        
        layout.addWidget(toolbar)
        
        # 文本编辑区
        self.text_edit = QTextEdit()
        self.text_edit.setPlaceholderText("✨ 开始编写文档...")
        self.text_edit.setFont(QFont("SF Pro Text", 12))
        self.text_edit.setMinimumHeight(500)
        layout.addWidget(self.text_edit)
        
        self.setLayout(layout)

    def new_file(self):
        self.text_edit.clear()
        self.current_file = None

    def open_file(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "打开 Word 文档", "", "Word Files (*.docx)")
        if file_path:
            try:
                doc = DocxDocument(file_path)
                full_text = [para.text for para in doc.paragraphs]
                self.text_edit.setText("\n".join(full_text))
                self.current_file = file_path
                QMessageBox.information(self, "成功", f"已打开文档: {os.path.basename(file_path)}")
            except Exception as e:
                QMessageBox.critical(self, "错误", f"无法打开文件：{str(e)}")

    def save_file(self):
        if self.current_file:
            self._save_to_file(self.current_file)
        else:
            self.save_as_file()

    def save_as_file(self):
        file_path, _ = QFileDialog.getSaveFileName(self, "保存 Word 文档", "", "Word Files (*.docx)")
        if file_path:
            self._save_to_file(file_path)
            self.current_file = file_path

    def _save_to_file(self, file_path):
        try:
            doc = DocxDocument()
            for line in self.text_edit.toPlainText().split('\n'):
                doc.add_paragraph(line)
            doc.save(file_path)
            QMessageBox.information(self, "成功", "文档已保存")
        except Exception as e:
            QMessageBox.critical(self, "错误", f"保存失败：{str(e)}")

    def bold(self):
        fmt = self.text_edit.textCursor().charFormat()
        fmt.setFontWeight(QFont.Bold if fmt.fontWeight() != QFont.Bold else QFont.Normal)
        self.text_edit.textCursor().mergeCharFormat(fmt)

    def italic(self):
        fmt = self.text_edit.textCursor().charFormat()
        fmt.setFontItalic(not fmt.fontItalic())
        self.text_edit.textCursor().mergeCharFormat(fmt)

# =============================== 表格编辑器 ===============================
class ExcelEditor(QWidget):
    def __init__(self):
        super().__init__()
        self.current_file = None
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout()
        layout.setContentsMargins(24, 24, 24, 24)
        layout.setSpacing(16)
        
        toolbar = QToolBar()
        toolbar.setMovable(False)
        
        new_action = QAction("📊 新建", self)
        new_action.triggered.connect(self.new_file)
        toolbar.addAction(new_action)
        
        open_action = QAction("📂 打开", self)
        open_action.triggered.connect(self.open_file)
        toolbar.addAction(open_action)
        
        save_action = QAction("💾 保存", self)
        save_action.triggered.connect(self.save_file)
        toolbar.addAction(save_action)
        
        saveas_action = QAction("📎 另存为", self)
        saveas_action.triggered.connect(self.save_as_file)
        toolbar.addAction(saveas_action)
        
        layout.addWidget(toolbar)

        self.table = QTableWidget(20, 10)
        self.table.setHorizontalHeaderLabels([f"列 {i+1}" for i in range(10)])
        self.table.setAlternatingRowColors(True)
        self.table.setStyleSheet("""
            QTableWidget::item:selected {
                background-color: #6366f1;
                color: white;
            }
            QTableWidget::item:hover {
                background-color: #f1f5f9;
            }
        """)
        layout.addWidget(self.table)
        self.setLayout(layout)

    def new_file(self):
        self.table.clear()
        self.table.setRowCount(20)
        self.table.setColumnCount(10)
        self.current_file = None

    def open_file(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "打开 Excel 文件", "", "Excel Files (*.xlsx *.xls)")
        if file_path:
            try:
                df = pd.read_excel(file_path)
                self.table.setRowCount(df.shape[0])
                self.table.setColumnCount(df.shape[1])
                self.table.setHorizontalHeaderLabels(df.columns.astype(str))
                for i in range(df.shape[0]):
                    for j in range(df.shape[1]):
                        item = QTableWidgetItem(str(df.iat[i, j]))
                        self.table.setItem(i, j, item)
                self.current_file = file_path
                QMessageBox.information(self, "成功", f"已打开表格: {os.path.basename(file_path)}")
            except Exception as e:
                QMessageBox.critical(self, "错误", f"无法打开文件：{str(e)}")

    def save_file(self):
        if self.current_file:
            self._save_to_file(self.current_file)
        else:
            self.save_as_file()

    def save_as_file(self):
        file_path, _ = QFileDialog.getSaveFileName(self, "保存 Excel 文件", "", "Excel Files (*.xlsx)")
        if file_path:
            self._save_to_file(file_path)
            self.current_file = file_path

    def _save_to_file(self, file_path):
        try:
            data = []
            for i in range(self.table.rowCount()):
                row = []
                for j in range(self.table.columnCount()):
                    item = self.table.item(i, j)
                    row.append(item.text() if item else "")
                data.append(row)
            df = pd.DataFrame(data)
            headers = [self.table.horizontalHeaderItem(j).text() for j in range(self.table.columnCount())]
            df.to_excel(file_path, index=False, header=headers)
            QMessageBox.information(self, "成功", "表格已保存")
        except Exception as e:
            QMessageBox.critical(self, "错误", f"保存失败：{str(e)}")

# =============================== PPT 编辑器与播放器 ===============================
class PPTEditor(QWidget):
    def __init__(self):
        super().__init__()
        self.presentation = None
        self.current_slide_index = 0
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout()
        layout.setContentsMargins(24, 24, 24, 24)
        layout.setSpacing(16)
        
        toolbar = QToolBar()
        toolbar.setMovable(False)
        
        new_action = QAction("🆕 新建", self)
        new_action.triggered.connect(self.new_ppt)
        toolbar.addAction(new_action)
        
        open_action = QAction("📂 打开", self)
        open_action.triggered.connect(self.open_ppt)
        toolbar.addAction(open_action)
        
        save_action = QAction("💾 保存", self)
        save_action.triggered.connect(self.save_ppt)
        toolbar.addAction(save_action)
        
        add_action = QAction("➕ 添加幻灯片", self)
        add_action.triggered.connect(self.add_slide)
        toolbar.addAction(add_action)
        
        prev_action = QAction("◀ 上一张", self)
        prev_action.triggered.connect(self.prev_slide)
        toolbar.addAction(prev_action)
        
        next_action = QAction("▶ 下一张", self)
        next_action.triggered.connect(self.next_slide)
        toolbar.addAction(next_action)
        
        play_action = QAction("🎬 播放", self)
        play_action.triggered.connect(self.play_slideshow)
        toolbar.addAction(play_action)
        
        layout.addWidget(toolbar)

        # 预览区域
        self.slide_label = QLabel("✨ 点击新建或打开 PPT 文件\n\n🎬 点击「播放」进入全屏放映模式\n\n💡 提示：全屏模式下支持鼠标/键盘翻页")
        self.slide_label.setAlignment(Qt.AlignCenter)
        self.slide_label.setStyleSheet("""
            background: rgba(255,255,255,0.9); 
            border-radius: 32px; 
            padding: 60px; 
            font-size: 14px;
            color: #64748b;
        """)
        self.slide_label.setMinimumHeight(450)
        layout.addWidget(self.slide_label)
        self.setLayout(layout)

    def new_ppt(self):
        self.presentation = Presentation()
        self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        self.current_slide_index = 0
        self.update_slide_display()
        QMessageBox.information(self, "成功", "已创建新的演示文稿")

    def open_ppt(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "打开 PPT 文件", "", "PowerPoint Files (*.pptx)")
        if file_path:
            try:
                self.presentation = Presentation(file_path)
                self.current_slide_index = 0
                self.update_slide_display()
                QMessageBox.information(self, "成功", f"已打开演示文稿: {os.path.basename(file_path)}")
            except Exception as e:
                QMessageBox.critical(self, "错误", f"无法打开文件：{str(e)}")

    def save_ppt(self):
        if not self.presentation:
            QMessageBox.warning(self, "警告", "没有打开的PPT")
            return
        file_path, _ = QFileDialog.getSaveFileName(self, "保存 PPT 文件", "", "PowerPoint Files (*.pptx)")
        if file_path:
            try:
                self.presentation.save(file_path)
                QMessageBox.information(self, "成功", "PPT已保存")
            except Exception as e:
                QMessageBox.critical(self, "错误", f"保存失败：{str(e)}")

    def add_slide(self):
        if self.presentation:
            self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
            self.current_slide_index = len(self.presentation.slides) - 1
            self.update_slide_display()

    def prev_slide(self):
        if self.presentation and self.current_slide_index > 0:
            self.current_slide_index -= 1
            self.update_slide_display()

    def next_slide(self):
        if self.presentation and self.current_slide_index < len(self.presentation.slides) - 1:
            self.current_slide_index += 1
            self.update_slide_display()

    def update_slide_display(self):
        if self.presentation and len(self.presentation.slides) > 0:
            slide = self.presentation.slides[self.current_slide_index]
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            display_text = f"📄 幻灯片 {self.current_slide_index+1} / {len(self.presentation.slides)}\n\n" + "\n".join(texts)
            self.slide_label.setText(display_text)
        else:
            self.slide_label.setText("✨ 无幻灯片")

    def play_slideshow(self):
        """启动全屏放映"""
        if not self.presentation:
            QMessageBox.warning(self, "警告", "请先打开或新建一个 PPT")
            return
        self.slide_show = SlideShowWindow(self.presentation, self)
        self.slide_show.showFullScreen()
        self.slide_show.exec_()

# =============================== 视频编辑器 ===============================
class VideoEditor(QWidget):
    def __init__(self):
        super().__init__()
        self.current_video = None
        self.clip = None
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout()
        layout.setContentsMargins(24, 24, 24, 24)
        layout.setSpacing(16)

        toolbar = QToolBar()
        toolbar.setMovable(False)
        
        open_action = QAction("🎬 打开视频", self)
        open_action.triggered.connect(self.open_video)
        toolbar.addAction(open_action)
        
        export_action = QAction("✂️ 导出剪辑", self)
        export_action.triggered.connect(self.export_clip)
        toolbar.addAction(export_action)
        
        layout.addWidget(toolbar)

        # 视频预览
        self.video_widget = QVideoWidget()
        self.video_widget.setStyleSheet("border-radius: 28px; background: black;")
        self.video_widget.setMinimumHeight(400)
        self.media_player = QMediaPlayer()
        self.media_player.setVideoOutput(self.video_widget)
        layout.addWidget(self.video_widget, 3)

        # 控制栏
        controls = QHBoxLayout()
        controls.setSpacing(12)
        self.play_button = QPushButton("▶ 播放")
        self.play_button.clicked.connect(self.toggle_play)
        self.play_button.setMinimumWidth(80)
        self.time_slider = QSlider(Qt.Horizontal)
        self.time_slider.setRange(0, 1000)
        self.time_slider.sliderMoved.connect(self.set_position)
        controls.addWidget(self.play_button)
        controls.addWidget(self.time_slider)
        layout.addLayout(controls)

        # 剪辑参数
        clip_group = QGroupBox("✂️ 剪辑设置")
        clip_layout = QGridLayout()
        clip_layout.setSpacing(12)
        clip_layout.addWidget(QLabel("开始时间 (秒):"), 0, 0)
        self.start_spin = QSpinBox()
        self.start_spin.setRange(0, 999999)
        self.start_spin.setStyleSheet("min-width: 100px;")
        clip_layout.addWidget(self.start_spin, 0, 1)
        clip_layout.addWidget(QLabel("结束时间 (秒):"), 1, 0)
        self.end_spin = QSpinBox()
        self.end_spin.setRange(0, 999999)
        self.end_spin.setStyleSheet("min-width: 100px;")
        clip_layout.addWidget(self.end_spin, 1, 1)
        clip_group.setLayout(clip_layout)
        layout.addWidget(clip_group)

        self.setLayout(layout)

        self.media_player.positionChanged.connect(self.position_changed)
        self.media_player.durationChanged.connect(self.duration_changed)

    def open_video(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "打开视频文件", "", "Video Files (*.mp4 *.avi *.mov)")
        if file_path:
            self.current_video = file_path
            self.media_player.setMedia(QMediaContent(QUrl.fromLocalFile(file_path)))
            self.media_player.play()
            self.play_button.setText("⏸ 暂停")
            try:
                self.clip = VideoFileClip(file_path)
                duration = self.clip.duration
                self.start_spin.setMaximum(int(duration))
                self.end_spin.setMaximum(int(duration))
                self.end_spin.setValue(int(duration))
                QMessageBox.information(self, "成功", f"已打开视频: {os.path.basename(file_path)}")
            except Exception as e:
                QMessageBox.warning(self, "警告", f"无法获取视频信息：{str(e)}")

    def toggle_play(self):
        if self.media_player.state() == QMediaPlayer.PlayingState:
            self.media_player.pause()
            self.play_button.setText("▶ 播放")
        else:
            self.media_player.play()
            self.play_button.setText("⏸ 暂停")

    def set_position(self, position):
        self.media_player.setPosition(position)

    def position_changed(self, position):
        self.time_slider.blockSignals(True)
        self.time_slider.setValue(position)
        self.time_slider.blockSignals(False)

    def duration_changed(self, duration):
        self.time_slider.setRange(0, duration)

    def export_clip(self):
        if not self.clip:
            QMessageBox.warning(self, "警告", "请先打开视频")
            return
        start = self.start_spin.value()
        end = self.end_spin.value()
        if start >= end:
            QMessageBox.critical(self, "错误", "开始时间必须小于结束时间")
            return
        output_path, _ = QFileDialog.getSaveFileName(self, "导出剪辑", "", "MP4 Files (*.mp4)")
        if output_path:
            try:
                subclip = self.clip.subclip(start, end)
                subclip.write_videofile(output_path, codec='libx264', audio_codec='aac', verbose=False, logger=None)
                QMessageBox.information(self, "成功", f"剪辑已导出至 {output_path}")
            except Exception as e:
                QMessageBox.critical(self, "错误", f"导出失败：{str(e)}")

# =============================== 主窗口 ===============================
class InspirationStudio(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("✨ Inspiration Studio 全能办公套件 | 创意无限 · 灵感随行")
        self.setGeometry(100, 100, 1400, 900)
        self.setStyleSheet(STYLE_SHEET)

        # 主窗口阴影
        add_shadow(self, blur_radius=40, offset=(0, 8), alpha=40)

        # 中心选项卡
        self.tabs = QTabWidget()
        self.tabs.setDocumentMode(True)
        self.setCentralWidget(self.tabs)

        # 创建各模块
        self.word_editor = WordEditor()
        self.excel_editor = ExcelEditor()
        self.ppt_editor = PPTEditor()
        self.video_editor = VideoEditor()

        self.tabs.addTab(self.word_editor, "📄 文字处理")
        self.tabs.addTab(self.excel_editor, "📊 表格处理")
        self.tabs.addTab(self.ppt_editor, "📽️ 幻灯片制作")
        self.tabs.addTab(self.video_editor, "🎬 视频剪辑")

        # 状态栏
        self.status_bar = QStatusBar()
        self.setStatusBar(self.status_bar)
        self.status_bar.showMessage("✨ Inspiration Studio 已就绪 | 让创意自由流淌")

        # 菜单栏
        self.create_menus()

    def create_menus(self):
        menubar = self.menuBar()
        
        # 文件菜单
        file_menu = menubar.addMenu("文件")
        exit_action = QAction("退出", self)
        exit_action.setShortcut("Ctrl+Q")
        exit_action.triggered.connect(self.close)
        file_menu.addAction(exit_action)
        
        # 帮助菜单
        help_menu = menubar.addMenu("帮助")
        about_action = QAction("关于 Inspiration Studio", self)
        about_action.triggered.connect(self.show_about)
        help_menu.addAction(about_action)

    def show_about(self):
        QMessageBox.about(self, "关于 Inspiration Studio",
                          "✨ Inspiration Studio 全能办公套件\n\n"
                          "版本 4.0 | 液态玻璃设计 · 全屏放映\n"
                          "基于 Python + PyQt5 构建\n"
                          "包含文档、表格、幻灯片及视频剪辑功能\n\n"
                          "© 2026 DreamWorks Studio\n"
                          "让创意无限 · 让灵感随行")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    
    # 设置应用字体
    font = QFont("Segoe UI", 10)
    app.setFont(font)
    
    # 创建主窗口
    window = InspirationStudio()
    window.show()
    
    sys.exit(app.exec_())